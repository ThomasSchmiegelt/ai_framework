"""
Dokument-Export für AI_Framework_Thomas — erzeugt DOCX, XLSX und PPTX aus Chat-/Canvas-Inhalten.

Alle Exporte tragen die AI_Framework_Thomas-Fußzeile ("KI generierter Inhalt" + optional
Name/Firma aus dem Nutzerprofil). Branding-Bilder (Vorlagen-Deckblatt,
Vorlagen-Kopfzeile, Logo) kommen aus dem Nutzerprofil (``data/profile_assets/``)
— sind keine hinterlegt, werden die Folien/Dokumente schlicht ohne Bild erzeugt.
DOCX kann auf Wunsch eine Kopfzeile als Bild einbetten (Flag
``_include_header_image``). KI-Antworten werden im DOCX mit "▶ Von KI generiert"
gekennzeichnet.

Einstiegspunkte: :func:`to_docx`, :func:`to_xlsx`, :func:`to_pptx`
(jede gibt einen Pfad zur temporären Ausgabedatei zurück).
"""
import re
import tempfile
import time
from pathlib import Path

# Inline-Markdown: ***fett-kursiv***, **fett**, *kursiv*, `code`
_MD_INLINE = re.compile(r"\*\*\*(.+?)\*\*\*|\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`", re.DOTALL)


def _strip_md(text: str) -> str:
    """Entfernt Inline-Markdown-Marker (für Überschriften/PDF, wo kein Mischsatz)."""
    return _MD_INLINE.sub(lambda m: next(g for g in m.groups() if g is not None), text)


def _add_md_runs(paragraph, text: str) -> None:
    """Fügt Text mit Inline-Markdown als echte Word-Runs hinzu (fett/kursiv/Code),
    damit in Word keine literalen ** / *** / ` erscheinen."""
    pos = 0
    for m in _MD_INLINE.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        bi, b, i, code = m.groups()
        if bi is not None:
            r = paragraph.add_run(bi); r.bold = True; r.italic = True
        elif b is not None:
            r = paragraph.add_run(b); r.bold = True
        elif i is not None:
            r = paragraph.add_run(i); r.italic = True
        elif code is not None:
            r = paragraph.add_run(code); r.font.name = "Consolas"
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


# ── Mathe-/LaTeX-Formeln ───────────────────────────────
# Formel-Segmente in vier Schreibweisen: $$..$$, $..$, \[..\], \(..\)
_MATH_RE = re.compile(r"\$\$.+?\$\$|\$[^$\n]+?\$|\\\[.+?\\\]|\\\(.+?\\\)", re.DOTALL)
_SENT = "\x00"          # geschuetztes Leerzeichen innerhalb einer Formel (Umbruch)
_MATHTEXT_PARSER = None


def _has_math(text: str) -> bool:
    return bool(_MATH_RE.search(text or ""))


def _math_to_mpl(text: str) -> str:
    """Vereinheitlicht alle Formel-Trenner auf ``$..$`` -- matplotlib-mathtext
    kennt nur das Dollar-Zeichen (kein Backslash-Klammer / $$)."""
    text = re.sub(r"\$\$(.+?)\$\$", r"$\1$", text, flags=re.DOTALL)
    text = re.sub(r"\\\[(.+?)\\\]", r"$\1$", text, flags=re.DOTALL)
    text = re.sub(r"\\\((.+?)\\\)", r"$\1$", text, flags=re.DOTALL)
    return text


def _mpl_safe_math(text: str) -> str:
    """Macht eine (bereits auf ``$..$`` normalisierte) Zeile zeichensicher fuer
    matplotlib: jede Formel wird mit dem mathtext-Parser geprueft -- nicht
    renderbare Formeln fallen auf ihren reinen Innentext zurueck; einzelne
    ungepaarte ``$`` werden maskiert. So stuerzt ``savefig`` nie ab."""
    global _MATHTEXT_PARSER
    if _MATHTEXT_PARSER is None:
        from matplotlib import mathtext
        _MATHTEXT_PARSER = mathtext.MathTextParser("agg")
    out, pos = [], 0
    for m in re.finditer(r"\$[^$\n]+?\$", text):
        out.append(text[pos:m.start()].replace("$", r"\$"))
        seg = m.group(0)
        try:
            _MATHTEXT_PARSER.parse(seg)
            out.append(seg)
        except Exception:
            out.append(seg.strip("$"))
        pos = m.end()
    out.append(text[pos:].replace("$", r"\$"))
    return "".join(out)


def _strip_md_keep_math(text: str) -> str:
    """Wie :func:`_strip_md`, laesst aber ``$..$``-Formeln unangetastet
    (sonst wuerde ``*`` in einer Formel als Kursiv-Marker verschluckt)."""
    saved: list[str] = []

    def _keep(m):
        saved.append(m.group(0))
        return "\x01%d\x02" % (len(saved) - 1)

    t = _MATH_RE.sub(_keep, text)
    t = _strip_md(t)
    for i, s in enumerate(saved):
        t = t.replace("\x01%d\x02" % i, s)
    return t


def _wrap_keep_math(text: str, width: int) -> list:
    """Zeilenumbruch wie :func:`textwrap.wrap`, bricht aber nie innerhalb einer
    Formel (Leerzeichen darin werden vor dem Umbruch geschuetzt)."""
    import textwrap
    protected = _MATH_RE.sub(lambda m: m.group(0).replace(" ", _SENT), text)
    lines = textwrap.wrap(protected, width=width,
                          break_long_words=False, break_on_hyphens=False) or [""]
    return [ln.replace(_SENT, " ") for ln in lines]


# Branding-Bilder aus dem Nutzerprofil (vom Nutzer im Profil hochgeladen)
_ASSETS_DIR   = Path(__file__).parent.parent / "data" / "profile_assets"
_DECKBLATT    = _ASSETS_DIR / "cover.jpg"     # Vorlagen-Deckblatt (Titelfolie)
_KOPFZEILE    = _ASSETS_DIR / "header.jpg"    # Vorlagen-Kopfzeile (Banner)

# Corporate-Farben (Primärpalette)
_CORP_PRIMARY   = "3B76BA"   # rgb(59,118,186)
_CORP_DARK      = "11314F"   # rgb(17,49,79)
_CORP_DEEP      = "003A74"   # rgb(0,58,116)
_CORP_MED       = "0056AD"   # rgb(0,86,173)
_CORP_LIGHT_TXT = "D4E8F8"   # rgb(212,232,248)
_CORP_DIM_TXT   = "A3C8EB"   # rgb(163,200,235)
_CORP_GRAY      = "6C6F76"   # rgb(108,111,118)


def _footer_text(data: dict) -> str:
    parts = []
    profile = data.get("_profile", {}) or {}
    name = " ".join(filter(None, [profile.get("first_name"), profile.get("last_name")])).strip()
    company = profile.get("company", "").strip()
    if name:
        parts.append(name)
    if company:
        parts.append(company)
    parts.append("KI generierter Inhalt")
    return " · ".join(parts)


# ── DOCX ─────────────────────────────────────────────────────────────────────

def to_docx(data: dict) -> Path:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    if data.get("_include_header_image"):
        _DOKUMENTE = _ASSETS_DIR / "header.jpg"   # Vorlagen-Kopfzeile aus dem Profil
        if _DOKUMENTE.exists():
            from docx.shared import Inches
            doc.add_picture(str(_DOKUMENTE), width=Inches(6.3))
            doc.add_paragraph()  # Abstand

        note_para = doc.add_paragraph()
        note_run = note_para.add_run("Dieser Bericht wurde von KI (AI_Framework_Thomas) generiert · " + time.strftime("%d.%m.%Y"))
        note_run.font.size = Pt(9)
        note_run.font.italic = True
        note_run.font.color.rgb = RGBColor(0x6C, 0x6F, 0x76)

    title = data.get("title", "Dokument")
    doc.add_heading(title, 0)

    content = data.get("content", "")
    if content:
        for para in content.split("\n"):
            if para.startswith("#### "):
                doc.add_heading(_strip_md(para[5:]), 4)
            elif para.startswith("### "):
                doc.add_heading(_strip_md(para[4:]), 3)
            elif para.startswith("## "):
                doc.add_heading(_strip_md(para[3:]), 2)
            elif para.startswith("# "):
                doc.add_heading(_strip_md(para[2:]), 1)
            elif para.startswith("- ") or para.startswith("* "):
                _add_md_runs(doc.add_paragraph(style="List Bullet"), para[2:])
            elif re.match(r"^\d+\.\s", para):
                _add_md_runs(doc.add_paragraph(style="List Number"),
                             re.sub(r"^\d+\.\s", "", para))
            elif para.strip():
                _add_md_runs(doc.add_paragraph(), para)

    messages = data.get("messages", [])
    for msg in messages:
        role = msg.get("role", "")
        text = msg.get("content", "")
        if role == "user":
            p = doc.add_paragraph()
            p.add_run("Benutzer: ").bold = True
            _add_md_runs(p, text)
        elif role == "assistant":
            ki_label = doc.add_paragraph()
            ki_run = ki_label.add_run("▶ Von KI generiert")
            ki_run.font.size = Pt(7)
            ki_run.font.italic = True
            ki_run.font.color.rgb = RGBColor(0x3B, 0x76, 0xBA)
            p = doc.add_paragraph()
            p.add_run("Assistent: ").bold = True
            _add_md_runs(p, text)

    headers = data.get("headers")
    rows    = data.get("rows")
    if headers and rows:
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"
        for i, h in enumerate(headers):
            table.cell(0, i).text = str(h)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                if ci < len(headers):
                    table.cell(ri + 1, ci).text = str(val) if val is not None else ""

    # Fußzeile: NUR Text, kein Bild (per Nutzeranweisung)
    footer_text = _footer_text(data)
    section = doc.sections[0]
    footer  = section.footer
    fp_para = footer.paragraphs[0]
    fp_para.text = footer_text
    fp_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if fp_para.runs:
        fp_para.runs[0].font.size = Pt(8)
        fp_para.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    fp = Path(tempfile.mktemp(suffix=".docx"))
    doc.save(str(fp))
    return fp


# ── XLSX ─────────────────────────────────────────────────────────────────────

def to_xlsx(data: dict) -> Path:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = data.get("title", "Tabelle")[:31]

    headers = data.get("headers", [])
    rows    = data.get("rows",    [])

    header_fill = PatternFill(start_color=_CORP_PRIMARY, end_color=_CORP_PRIMARY, fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    for ci, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=ci, value=str(h))
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    for ri, row in enumerate(rows, 2):
        for ci, val in enumerate(row, 1):
            if ci <= len(headers):
                ws.cell(row=ri, column=ci, value=val)

    for col in ws.columns:
        max_len = max((len(str(cell.value or "")) for cell in col), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    footer_text = _footer_text(data)
    ws.oddFooter.center.text = footer_text
    ws.oddFooter.center.size = 8

    fp = Path(tempfile.mktemp(suffix=".xlsx"))
    wb.save(str(fp))
    return fp


# ── PPTX ─────────────────────────────────────────────────────────────────────

def to_pptx(data: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # Kopfzeilen-Höhe in Inches (proportional zum Bild 90/755 × 7.5")
    _HEADER_H = Inches(0.89)
    _SLIDE_W  = prs.slide_width
    _SLIDE_H  = prs.slide_height

    # Corporate-Farben als RGBColor
    C_DARK      = RGBColor(0x11, 0x31, 0x4F)
    C_DEEP      = RGBColor(0x00, 0x3A, 0x74)
    C_PRIMARY   = RGBColor(0x3B, 0x76, 0xBA)
    C_MED       = RGBColor(0x00, 0x56, 0xAD)
    C_LIGHT_TXT = RGBColor(0xD4, 0xE8, 0xF8)
    C_DIM_TXT   = RGBColor(0xA3, 0xC8, 0xEB)
    C_GRAY      = RGBColor(0x6C, 0x6F, 0x76)
    C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

    slides = data.get("slides", [])
    footer_text = _footer_text(data)

    for slide_idx, slide_data in enumerate(slides):
        slide  = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "bullets")
        title  = slide_data.get("title", "")

        # ── Deck­blatt (erste Folie oder explizit "title") ────────────────────
        if slide_idx == 0 and layout == "title" and _DECKBLATT.exists():
            # Deckblatt als vollflächiger Hintergrund
            slide.shapes.add_picture(
                str(_DECKBLATT), 0, 0, _SLIDE_W, _SLIDE_H
            )
            # Titel im blauen Bereich (untere Hälfte der Hexagon-Fläche)
            _add_text_run(
                slide, title,
                l=Inches(0.6), t=Inches(2.6), w=Inches(7.5), h=Inches(1.4),
                size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
            )
            content = slide_data.get("content", "")
            if content:
                _add_text_run(
                    slide, content,
                    l=Inches(0.6), t=Inches(4.2), w=Inches(7.5), h=Inches(0.8),
                    size=Pt(20), bold=False, color=C_DIM_TXT, align=PP_ALIGN.LEFT
                )

        # ── Inhaltsfolie (alle anderen) ───────────────────────────────────────
        else:
            # Weißer Hintergrund
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = C_WHITE

            # Kopfzeile als Bild-Strip
            if _KOPFZEILE.exists():
                slide.shapes.add_picture(
                    str(_KOPFZEILE), 0, 0, _SLIDE_W, _HEADER_H
                )

            # Dünner Corporate-Trennstrich unterhalb der Kopfzeile
            _content_top = _HEADER_H + Inches(0.05)

            if layout == "section":
                # Abschnitts-Folie: Corporate-blauer Block
                from pptx.util import Emu
                section_box = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    0, _content_top, _SLIDE_W, _SLIDE_H - _content_top
                )
                section_box.fill.solid()
                section_box.fill.fore_color.rgb = C_DEEP
                section_box.line.fill.background()
                _add_text_run(
                    slide, title,
                    l=Inches(1), t=_content_top + Inches(1.8), w=Inches(11.33), h=Inches(1.5),
                    size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER
                )
                subtitle = slide_data.get("subtitle") or slide_data.get("content", "")
                if subtitle:
                    _add_text_run(
                        slide, subtitle,
                        l=Inches(1.5), t=_content_top + Inches(3.4), w=Inches(10.33), h=Inches(0.8),
                        size=Pt(22), bold=False, color=C_DIM_TXT, align=PP_ALIGN.CENTER
                    )

            else:
                # Folientitel (dunkler Streifen in Corporate-Blau)
                from pptx.util import Emu
                title_bg = slide.shapes.add_shape(
                    1,
                    0, _content_top, _SLIDE_W, Inches(0.85)
                )
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = C_DARK
                title_bg.line.fill.background()

                if title:
                    _add_text_run(
                        slide, title,
                        l=Inches(0.4), t=_content_top + Inches(0.1), w=Inches(12.5), h=Inches(0.7),
                        size=Pt(26), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
                    )

                body_top = _content_top + Inches(1.0)
                body_h   = _SLIDE_H - body_top - Inches(0.5)

                bullets = slide_data.get("bullets", [])
                if isinstance(bullets, str):
                    bullets = [b for b in bullets.split("\n") if b.strip()]

                left  = slide_data.get("left")
                right = slide_data.get("right")
                img_r = slide_data.get("image_right")

                if left and (right or img_r):
                    # Zweispaltig
                    _add_bullets(slide, _to_lines(left),
                                 Inches(0.4), body_top, Inches(5.9), body_h,
                                 C_DARK, C_GRAY)
                    if img_r:
                        import base64, io as _io
                        _embed_b64_image(slide, img_r,
                                         Inches(7.1), body_top, Inches(5.8), body_h)
                    elif right:
                        _add_bullets(slide, _to_lines(right),
                                     Inches(7.1), body_top, Inches(5.9), body_h,
                                     C_DEEP, C_GRAY)
                elif bullets:
                    _add_bullets(slide, bullets,
                                 Inches(0.6), body_top, Inches(12.13), body_h,
                                 C_DARK, C_GRAY)
                else:
                    content = slide_data.get("content", "")
                    if content:
                        _add_text_run(
                            slide, content,
                            l=Inches(0.6), t=body_top, w=Inches(12.13), h=body_h,
                            size=Pt(20), bold=False, color=C_DARK, align=PP_ALIGN.LEFT
                        )

        # Fußzeile auf jeder Folie
        ftBox = slide.shapes.add_textbox(
            Inches(0), _SLIDE_H - Inches(0.3), _SLIDE_W, Inches(0.3)
        )
        ftf = ftBox.text_frame
        fp_p = ftf.paragraphs[0]
        fp_p.text = footer_text
        fp_p.alignment = PP_ALIGN.CENTER
        if fp_p.runs:
            fp_p.runs[0].font.size = Pt(8)
            fp_p.runs[0].font.color.rgb = C_GRAY

    fp = Path(tempfile.mktemp(suffix=".pptx"))
    prs.save(str(fp))
    return fp


# ── Hilfsfunktionen ──────────────────────────────────────────────────────────

def _add_text_run(slide, text: str, l, t, w, h, size, bold: bool, color, align):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color


def _add_bullets(slide, items: list, l, t, w, h, text_color, bullet_color):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"◆  {item}" if not str(item).startswith(("•", "◆", "-", "*")) else item
        run.font.size  = Pt(19)
        run.font.color.rgb = text_color


def _embed_b64_image(slide, data_url: str, l, t, w, h):
    """Bettet ein Base64-Bild seitenverhältnis-erhaltend (contain) in die Box
    (l, t, w, h) ein und zentriert es darin – verhindert Verzerrung."""
    import base64, io
    try:
        header, b64 = data_url.split(",", 1)
        img_bytes = base64.b64decode(b64)
        iw = ih = 0
        try:
            from PIL import Image
            with Image.open(io.BytesIO(img_bytes)) as im:
                iw, ih = im.size
        except Exception:
            iw = ih = 0
        if iw > 0 and ih > 0:
            from pptx.util import Emu
            box_w, box_h = float(w), float(h)
            aspect = iw / ih
            draw_w = box_w
            draw_h = draw_w / aspect
            if draw_h > box_h:
                draw_h = box_h
                draw_w = draw_h * aspect
            off_l = float(l) + (box_w - draw_w) / 2
            off_t = float(t) + (box_h - draw_h) / 2
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Emu(int(off_l)), Emu(int(off_t)), Emu(int(draw_w)), Emu(int(draw_h)),
            )
        else:
            slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)
    except Exception:
        pass


def _to_lines(val) -> list:
    if isinstance(val, list):
        return [str(b) for b in val]
    return [b for b in str(val).split("\n")]


# ── PDF (matplotlib, ohne LaTeX/TeX-Installation) ────────────────────────────

# Corporate-Farben als matplotlib-RGB-Tupel
_PDF_DARK   = (0x11 / 255, 0x31 / 255, 0x4F / 255)
_PDF_ACCENT = (0x3B / 255, 0x76 / 255, 0xBA / 255)
_PDF_GRAY   = (0x6C / 255, 0x6F / 255, 0x76 / 255)
_PDF_WHITE  = (1, 1, 1)


def to_pdf(data: dict) -> Path:
    """Erzeugt ein PDF aus einem Dokument ({title, content}) ODER einer
    Präsentation ({type:"presentation", slides:[…]}). Reines matplotlib —
    keine TeX-/LaTeX-Installation nötig."""
    import matplotlib
    matplotlib.use("Agg")
    fp = Path(tempfile.mktemp(suffix=".pdf"))
    if data.get("type") == "presentation":
        _pdf_presentation(data, fp)
    else:
        _pdf_document(data, fp)
    return fp


def _pdf_text(ax, x, y, s, **kw):
    """``ax.text`` mit automatischer Formel-Erkennung: enthält ``s`` eine
    ``$..$``-Formel, wird sie via mathtext gesetzt (vorher entschärft, damit
    ``savefig`` nie abstürzt), sonst als reiner Text."""
    if _has_math(s):
        ax.text(x, y, _mpl_safe_math(s), parse_math=True, **kw)
    else:
        ax.text(x, y, s, parse_math=False, **kw)


def _pdf_blocks_from_markdown(content: str) -> list[dict]:
    """Markdown → einfache Blockliste {text,size,bold,bullet,gap}. Inline-Marker
    werden entfernt, ``$..$``-Formeln aber erhalten (matplotlib-mathtext)."""
    def _clean(s: str) -> str:
        return _strip_md_keep_math(_math_to_mpl(s))

    blocks: list[dict] = []
    for raw in (content or "").split("\n"):
        line = raw.rstrip()
        if not line.strip():
            blocks.append({"text": "", "size": 5, "bold": False, "bullet": False, "gap": True})
            continue
        if line.startswith("#### "):
            blocks.append({"text": _clean(line[5:]), "size": 11, "bold": True, "bullet": False})
        elif line.startswith("### "):
            blocks.append({"text": _clean(line[4:]), "size": 12, "bold": True, "bullet": False})
        elif line.startswith("## "):
            blocks.append({"text": _clean(line[3:]), "size": 13.5, "bold": True, "bullet": False})
        elif line.startswith("# "):
            blocks.append({"text": _clean(line[2:]), "size": 15, "bold": True, "bullet": False})
        elif line.startswith("- ") or line.startswith("* "):
            blocks.append({"text": _clean(line[2:]), "size": 10.5, "bold": False, "bullet": True})
        elif re.match(r"^\d+\.\s", line):
            blocks.append({"text": _clean(line), "size": 10.5, "bold": False, "bullet": False})
        else:
            blocks.append({"text": _clean(line), "size": 10.5, "bold": False, "bullet": False})
    return blocks


def _pdf_document(data: dict, fp: Path) -> None:
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    PW, PH = 8.27, 11.69          # A4 hochkant (Zoll)
    L, R, TOP, BOT = 0.10, 0.92, 0.93, 0.07
    usable_pt = (R - L) * PW * 72
    footer_text = _footer_text(data)
    title = data.get("title", "Dokument")

    blocks = _pdf_blocks_from_markdown(data.get("content", ""))

    pages: list = []
    fig = ax = None
    y = 0.0

    def _new_page():
        nonlocal fig, ax, y
        fig = plt.figure(figsize=(PW, PH))
        ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.set_autoscale_on(False)
        ax.text(0.5, BOT - 0.025, footer_text, ha="center", va="center",
                fontsize=8, color=_PDF_GRAY, parse_math=False)
        pages.append(fig)
        y = TOP

    def _advance(size, factor=1.55):
        nonlocal y
        y -= (size * factor) / (PH * 72)

    _new_page()
    # Titel
    _pdf_text(ax, L, y, _strip_md_keep_math(_math_to_mpl(title)), ha="left", va="top",
              fontsize=20, fontweight="bold", color=_PDF_DARK)
    _advance(20, 2.2)
    ax.plot([L, R], [y, y], color=_PDF_ACCENT, lw=1.2)
    _advance(10)

    for blk in blocks:
        if blk.get("gap"):
            _advance(6)
            continue
        size = blk["size"]
        chars = max(20, int(usable_pt / (0.50 * size)))
        prefix = "•  " if blk["bullet"] else ""
        indent = 0.022 if blk["bullet"] else 0.0
        wrapped = _wrap_keep_math(blk["text"], chars)
        for i, ln in enumerate(wrapped):
            if y < BOT + 0.02:
                _new_page()
            x = L + (indent if (blk["bullet"] and i > 0) else 0)
            _pdf_text(ax, x, y, (prefix + ln) if (blk["bullet"] and i == 0) else ln,
                      ha="left", va="top", fontsize=size,
                      fontweight="bold" if blk["bold"] else "normal",
                      color=_PDF_DARK if blk["bold"] else (0.1, 0.1, 0.1))
            _advance(size)
        if blk["bold"]:
            _advance(4)

    with PdfPages(str(fp)) as pdf:
        for f in pages:
            pdf.savefig(f)
            plt.close(f)


def _slide_bullets(sd: dict) -> list:
    """Sammelt die Aufzählungspunkte einer Folie aus ``bullets`` / ``left`` +
    ``right`` / ``content`` (für PDF- und LaTeX-Export gleichermaßen)."""
    bullets = sd.get("bullets")
    if isinstance(bullets, str):
        bullets = [b for b in bullets.split("\n") if b.strip()]
    if not bullets:
        left, right = sd.get("left"), sd.get("right")
        if left or right:
            bullets = _to_lines(left or "") + _to_lines(right or "")
        elif sd.get("content"):
            bullets = [b for b in str(sd["content"]).split("\n") if b.strip()]
    return bullets or []


def _pdf_presentation(data: dict, fp: Path) -> None:
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    from matplotlib.patches import Rectangle

    SW, SH = 13.33, 7.5           # 16:9 quer (Zoll)
    footer_text = _footer_text(data)
    slides = data.get("slides", []) or []

    with PdfPages(str(fp)) as pdf:
        for sd in slides:
            layout = sd.get("layout", "bullets")
            title = sd.get("title", "")
            fig = plt.figure(figsize=(SW, SH))
            ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
            ax.set_xlim(0, 1); ax.set_ylim(0, 1)

            if layout == "title":
                ax.add_patch(Rectangle((0, 0), 1, 1, color=_PDF_DARK))
                ax.text(0.5, 0.56, title, ha="center", va="center", fontsize=34,
                        fontweight="bold", color=_PDF_WHITE, wrap=True, parse_math=False)
                sub = sd.get("content", "")
                if sub:
                    ax.text(0.5, 0.40, sub, ha="center", va="center", fontsize=18,
                            color=(0.64, 0.78, 0.92), wrap=True, parse_math=False)
            elif layout == "section":
                ax.add_patch(Rectangle((0, 0), 1, 1, color=(0, 0x3A / 255, 0x74 / 255)))
                ax.text(0.5, 0.55, title, ha="center", va="center", fontsize=30,
                        fontweight="bold", color=_PDF_WHITE, wrap=True, parse_math=False)
                sub = sd.get("subtitle") or sd.get("content", "")
                if sub:
                    ax.text(0.5, 0.40, sub, ha="center", va="center", fontsize=18,
                            color=(0.64, 0.78, 0.92), wrap=True, parse_math=False)
            else:
                # Titelstreifen
                ax.add_patch(Rectangle((0, 0.86), 1, 0.14, color=_PDF_DARK))
                ax.text(0.04, 0.93, title, ha="left", va="center", fontsize=22,
                        fontweight="bold", color=_PDF_WHITE, parse_math=False)
                # Inhalt sammeln (Formeln bleiben erhalten)
                bullets = _slide_bullets(sd)
                yb = 0.78
                for b in bullets:
                    if yb < 0.08:
                        break
                    clean = _strip_md_keep_math(_math_to_mpl(str(b)))
                    for i, ln in enumerate(_wrap_keep_math(clean, 78)):
                        txt = ("◆  " + ln) if i == 0 else "    " + ln
                        _pdf_text(ax, 0.05, yb, txt, ha="left", va="top",
                                  fontsize=17, color=_PDF_DARK)
                        yb -= 0.062

            ax.text(0.5, 0.02, footer_text, ha="center", va="center",
                    fontsize=8, color=_PDF_GRAY, parse_math=False)
            pdf.savefig(fig)
            plt.close(fig)


# == LaTeX-Quellcode (reine .tex-Datei) =====================================
# Erzeugt echtes LaTeX: Dokument -> article, Präsentation -> beamer. Formeln
# ($..$, $$..$$, \(..\), \[..\]) bleiben als echtes LaTeX-Math erhalten;
# Markdown (**fett**, *kursiv*, `code`, #/##/### Überschriften, -/* und 1.
# Listen) wird in LaTeX-Befehle übersetzt. Reiner Text wird LaTeX-sicher
# maskiert -- aber niemals innerhalb von Formeln.

def _latex_escape_text(s: str) -> str:
    """Maskiert LaTeX-Sonderzeichen in reinem Text (nicht in Formeln)."""
    s = s.replace(chr(92), "\x00BS\x00")        # Backslash zuerst sichern
    for a, b in [("&", r"\&"), ("%", r"\%"), ("$", r"\$"), ("#", r"\#"),
                 ("_", r"\_"), ("{", r"\{"), ("}", r"\}"),
                 ("~", r"\textasciitilde{}"), ("^", r"\textasciicircum{}")]:
        s = s.replace(a, b)
    return s.replace("\x00BS\x00", r"\textbackslash{}")


def _math_to_latex(seg: str) -> str:
    """Normalisiert ein Formel-Segment auf gültiges LaTeX-Math: ``$$..$$`` ->
    ``\\[..\\]`` (display); ``$..$``, ``\\(..\\)``, ``\\[..\\]`` bleiben gültig."""
    m = re.match(r"^\$\$(.+)\$\$$", seg, re.DOTALL)
    if m:
        return r"\[" + m.group(1) + r"\]"
    return seg


def _latex_inline(text: str) -> str:
    """Eine Textzeile -> LaTeX: Formeln werden geschützt, der Rest maskiert,
    Markdown-Inline (**fett**/*kursiv*/`code`) in LaTeX-Befehle übersetzt."""
    saved: list[str] = []

    def _keep(m):
        saved.append(m.group(0))
        return "\x01%d\x02" % (len(saved) - 1)

    t = _MATH_RE.sub(_keep, text)
    t = _latex_escape_text(t)
    t = re.sub(r"\*\*(.+?)\*\*", r"\\textbf{\1}", t)
    t = re.sub(r"\*(.+?)\*", r"\\textit{\1}", t)
    t = re.sub(r"`(.+?)`", r"\\texttt{\1}", t)
    for i, s in enumerate(saved):
        t = t.replace("\x01%d\x02" % i, _math_to_latex(s))
    return t


def _latex_table(rows: list[str]) -> str:
    """Markdown-Tabelle (Zeilen mit ``|``) -> LaTeX-tabular. Die Trennzeile
    (``|---|---|``) wird übersprungen."""
    def cells(line: str) -> list[str]:
        s = line.strip()
        if s.startswith("|"):
            s = s[1:]
        if s.endswith("|"):
            s = s[:-1]
        return [c.strip() for c in s.split("|")]

    body = [r for r in rows if not re.match(r"^\s*\|?[\s:|-]+\|?\s*$", r)]
    if not body:
        return ""
    ncol = max(len(cells(r)) for r in body)
    out = [r"\begin{center}", r"\begin{tabular}{|" + "l|" * ncol + "}", r"\hline"]
    for i, r in enumerate(body):
        cs = cells(r)
        cs += [""] * (ncol - len(cs))
        out.append(" & ".join(_latex_inline(c) for c in cs) + r" \\")
        out.append(r"\hline")
    out += [r"\end{tabular}", r"\end{center}"]
    return "\n".join(out)


def _md_to_latex_body(content: str) -> str:
    """Markdown -> LaTeX-Body. Blockorientiert: mehrzeilige Display-Formeln
    (``$$..$$`` / ``\\[..\\]``) und Codeblöcke (``` ``` ```) werden ZUERST aus
    dem Fließtext herausgelöst und unzerschnitten gesetzt — sonst zerbricht die
    zeilenweise Verarbeitung jede über mehrere Zeilen gehende Formel."""
    text = content or ""
    stash: list[str] = []

    def _keep_block(latex: str) -> str:
        stash.append(latex)
        return "\n\x03%d\x03\n" % (len(stash) - 1)

    # 1) Codeblöcke schützen (verbatim, NICHT maskiert)
    text = re.sub(
        r"```[A-Za-z0-9_+-]*\n(.*?)```",
        lambda m: _keep_block("\\begin{verbatim}\n" + m.group(1).rstrip("\n") + "\n\\end{verbatim}"),
        text, flags=re.DOTALL)
    # 2) Display-Formeln (auch mehrzeilig) schützen
    text = re.sub(
        r"\$\$(.+?)\$\$|\\\[(.+?)\\\]",
        lambda m: _keep_block("\\[\n" + (m.group(1) or m.group(2)).strip() + "\n\\]"),
        text, flags=re.DOTALL)

    out: list[str] = []
    list_mode = None  # "itemize" | "enumerate"

    def _close():
        nonlocal list_mode
        if list_mode:
            out.append(r"\end{%s}" % list_mode)
            list_mode = None

    lines = text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        stripped = line.strip()
        # Geschützter Block (Code/Formel) -> 1:1 einsetzen
        mblk = re.match(r"^\x03(\d+)\x03$", stripped)
        if mblk:
            _close(); out.append(stash[int(mblk.group(1))]); i += 1; continue
        # Markdown-Tabelle: aufeinanderfolgende Zeilen mit '|'
        if "|" in stripped and stripped.startswith("|"):
            tbl = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                tbl.append(lines[i]); i += 1
            _close(); out.append(_latex_table(tbl)); continue
        if not stripped:
            _close(); out.append("")
        elif line.startswith("#### "):
            _close(); out.append(r"\paragraph{%s}" % _latex_inline(line[5:]))
        elif line.startswith("### "):
            _close(); out.append(r"\subsubsection{%s}" % _latex_inline(line[4:]))
        elif line.startswith("## "):
            _close(); out.append(r"\subsection{%s}" % _latex_inline(line[3:]))
        elif line.startswith("# "):
            _close(); out.append(r"\section{%s}" % _latex_inline(line[2:]))
        elif line.startswith("- ") or line.startswith("* "):
            if list_mode != "itemize":
                _close(); out.append(r"\begin{itemize}"); list_mode = "itemize"
            out.append(r"  \item " + _latex_inline(line[2:]))
        elif re.match(r"^\d+\.\s", line):
            if list_mode != "enumerate":
                _close(); out.append(r"\begin{enumerate}"); list_mode = "enumerate"
            out.append(r"  \item " + _latex_inline(re.sub(r"^\d+\.\s+", "", line)))
        else:
            _close(); out.append(_latex_inline(line))
        i += 1
    _close()
    return "\n".join(out)


def _latex_document(data: dict) -> str:
    title = data.get("title", "Dokument")
    body = _md_to_latex_body(data.get("content", ""))
    footer = _footer_text(data)
    lang = (data.get("_profile", {}) or {}).get("lang", "de")
    babel = "english" if lang == "en" else "ngerman"
    pre = [
        r"\documentclass[11pt,a4paper]{article}",
        r"\usepackage[a4paper,margin=25mm]{geometry}",
        r"\usepackage[utf8]{inputenc}",
        r"\usepackage[T1]{fontenc}",
        r"\usepackage[%s]{babel}" % babel,
        r"\usepackage{amsmath,amssymb}",
        r"\usepackage{lmodern}",
        r"\usepackage{parskip}",
        r"\usepackage{hyperref}",
        r"\usepackage{fancyhdr}",
        r"\pagestyle{fancy}",
        r"\fancyhf{}",
        r"\cfoot{\thepage}",
        r"\lfoot{\small " + _latex_escape_text(footer) + "}",
        r"\setlength{\headheight}{14pt}",
        r"\title{" + _latex_inline(title) + "}",
        r"\date{\today}",
        r"\begin{document}",
        r"\maketitle",
    ]
    return "\n".join(pre) + "\n\n" + body + "\n\n" + r"\end{document}" + "\n"


def _latex_presentation(data: dict) -> str:
    title = data.get("title", "Präsentation")
    slides = data.get("slides", []) or []
    lang = (data.get("_profile", {}) or {}).get("lang", "de")
    babel = "english" if lang == "en" else "ngerman"
    out = [
        r"\documentclass[aspectratio=169]{beamer}",
        r"\usepackage[utf8]{inputenc}",
        r"\usepackage[T1]{fontenc}",
        r"\usepackage[%s]{babel}" % babel,
        r"\usepackage{amsmath,amssymb}",
        r"\usetheme{default}",
        r"\setbeamertemplate{navigation symbols}{}",
        r"\title{" + _latex_inline(title) + "}",
        r"\begin{document}",
        r"\frame{\titlepage}",
    ]
    for sd in slides:
        out.append(r"\begin{frame}{" + _latex_inline(sd.get("title", "")) + "}")
        bullets = _slide_bullets(sd)
        if bullets:
            out.append(r"\begin{itemize}")
            for b in bullets:
                out.append(r"  \item " + _latex_inline(str(b)))
            out.append(r"\end{itemize}")
        out.append(r"\end{frame}")
    out.append(r"\end{document}")
    return "\n".join(out) + "\n"


def to_latex(data: dict) -> Path:
    """Erzeugt eine reine LaTeX-Quelldatei (.tex) aus einem Dokument
    ({title, content}) ODER einer Präsentation ({type:"presentation", slides})."""
    if data.get("type") == "presentation":
        src = _latex_presentation(data)
    else:
        src = _latex_document(data)
    fp = Path(tempfile.mktemp(suffix=".tex"))
    fp.write_text(src, encoding="utf-8")
    return fp
